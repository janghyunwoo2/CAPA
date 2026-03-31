import sys
from pptx import Presentation

def extract_slide_content(pptx_file, slide_index):
    try:
        prs = Presentation(pptx_file)
        slide = prs.slides[slide_index]
        print(f"--- Slide {slide_index + 1} Content ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text.replace('\n', ' '))
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_slide_content(sys.argv[1], int(sys.argv[2]) - 1)
