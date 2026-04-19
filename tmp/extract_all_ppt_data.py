import sys
from pptx import Presentation

def extract_all_data(pptx_file):
    try:
        prs = Presentation(pptx_file)
        for i, slide in enumerate(prs.slides):
            print(f"=== Slide {i+1} ===")
            
            # Title
            title = "(No Title)"
            if slide.shapes.title:
                title = slide.shapes.title.text
            else:
                text_shapes = [s for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if text_shapes:
                    # Find the shape with the largest font size
                    def get_max_font_size(shape):
                        try:
                            return max(p.font.size or 0 for p in shape.text_frame.paragraphs)
                        except:
                            return 0
                    text_shapes.sort(key=get_max_font_size, reverse=True)
                    title = text_shapes[0].text.strip().replace('\n', ' ')
            print(f"TITLE: {title}")
            
            # Content
            print("CONTENT:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.replace('\n', ' ')
                    print(f"- {text}")
            
            # Notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                print(f"NOTES: {notes}")
            else:
                print("NOTES: (None)")
            print()
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_all_data(sys.argv[1])
