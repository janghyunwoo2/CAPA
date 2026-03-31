import sys
from pptx import Presentation

def extract_notes(pptx_file, slide_index=None):
    try:
        prs = Presentation(pptx_file)
        if slide_index is not None:
            if 0 <= slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    print(f"--- Slide {slide_index + 1} Notes ---")
                    print(notes)
                else:
                    print(f"Slide {slide_index + 1} has no notes.")
            else:
                print(f"Invalid slide index: {slide_index + 1}")
        else:
            for i, slide in enumerate(prs.slides):
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    print(f"--- Slide {i + 1} Notes ---")
                    print(notes)
                else:
                    print(f"Slide {i + 1} has no notes.")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_ppt_notes.py <file.pptx> [slide_number]")
    else:
        file_path = sys.argv[1]
        slide_num = int(sys.argv[2]) - 1 if len(sys.argv) > 2 else None
        extract_notes(file_path, slide_num)
