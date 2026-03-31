import sys
from pptx import Presentation

def extract_titles(pptx_file):
    try:
        prs = Presentation(pptx_file)
        for i, slide in enumerate(prs.slides):
            title = "(No Title)"
            text_shapes = [s for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if slide.shapes.title:
                title = slide.shapes.title.text
            elif text_shapes:
                # Find the shape with the largest font size in its first paragraph
                def get_max_font_size(shape):
                    try:
                        return max(p.font.size or 0 for p in shape.text_frame.paragraphs)
                    except:
                        return 0
                text_shapes.sort(key=get_max_font_size, reverse=True)
                title = text_shapes[0].text.strip().replace('\n', ' ')
            print(f"Slide {i+1}: {title}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_ppt_titles.py <file.pptx>")
    else:
        extract_titles(sys.argv[1])
